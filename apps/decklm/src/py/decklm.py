from pptx import Presentation
from pptx.util import Inches, Pt, Cm
import json
import os

def remove_trailing_period(s):
    """Removes the trailing period from a string if it exists."""
    if s.endswith('.'):
        return s[:-1]
    return s


def remove_body(slide):
    """Remove the body from a slide of layouts[1]"""
    for shape in slide.shapes:
        if shape.is_placeholder and shape.placeholder_format.idx == 1:  # Typically idx 1 is the body
            sp = shape._element  # Get the element to remove
            sp.getparent().remove(sp)

def make_pptx(file_path, json_data):
    """ 
    This function will generate the pptx document and save at path file_path.
    The data should include:
    1. An overall title for the document.
    2. Slides, and each slide has a title, with paragraphs and lists in any combination.
    """
    data = json.loads(json_data)
    bold_runs = []
    bullets = []
    bullets = []
    
    font_path = os.path.expanduser("~/common-lisp/ninx/apps/decklm/fonts/Roboto-Regular.ttf")
    font_family = "Roboto Medium"
    prs = Presentation()
    prs.slide_width = Inches(11.02)
    prs.slide_height = Inches(6.20)
    
    title_slide_layout = prs.slide_layouts[0]
    title_slide = prs.slides.add_slide(title_slide_layout)
    title = title_slide.shapes.title
    subtitle = title_slide.placeholders[1]
    subtitle.left = Cm(5)
    subtitle.top = Cm(10.8)
    title.text = data["overall-title"]
    title.top = Cm(4)
    title.left = Cm(3)
    
    outline_slide_layout = prs.slide_layouts[1]
    outline_slide = prs.slides.add_slide(outline_slide_layout)
    outline_title = outline_slide.shapes.title
    outline_title.text = "Outline"
    outline_title_para = outline_title.text_frame.paragraphs[0]
    outline_title_para.font.size = Pt(28)
    outline_title_para.font.bold = True
    outline_title.top = Cm(-0.5)
    outline_title.left = Cm(0)

    remove_body(outline_slide)
    outline_body = outline_slide.shapes.add_textbox(Cm(0.5), Cm(2), Inches(10.8), Inches(3.02))
    outline_text_frame = outline_body.text_frame

    for point in data["outline"]:
        p = outline_text_frame.add_paragraph()# if text_frame.text else text_frame
        p_run = p.add_run()
        bullets.append(p)
        p_run.text = point
        p_run.font.size = Pt(20)
    outline_text_frame.fit_text(font_family=font_family, font_file=font_path)

    definitions_slide_layout = prs.slide_layouts[1]
    definitions_slide = prs.slides.add_slide(definitions_slide_layout)
    definitions_title = definitions_slide.shapes.title
    definitions_title.text = "Definitions"
    definitions_title_para = definitions_title.text_frame.paragraphs[0]
    definitions_title_para.font.size = Pt(28)
    definitions_title_para.font.bold = True
    definitions_title.top = Cm(-0.5)
    definitions_title.left = Cm(0)
    
    remove_body(definitions_slide)
    definitions_body = definitions_slide.shapes.add_textbox(Cm(0.5), Cm(2), Inches(10.8), Inches(3.02))
    definitions_text_frame = definitions_body.text_frame

    for point in data["definitions"]:
        p = definitions_text_frame.add_paragraph()# if text_frame.text else text_frame
        bullets.append(p)
        p_run = p.add_run()
        term = remove_trailing_period(point["term"])
        definition = point["definition"]
        p_run.text = f"{term}: "
        p_run.font.bold = True
        bold_runs.append(p_run)
        p_run.font.size = Pt(20)
        e_run = p.add_run()
        e_run.text = definition
        e_run.font.size = Pt(20)
    definitions_text_frame.fit_text(font_family=font_family, font_file=font_path)
        
    for slide_data in data["content"]:
        slide_layout = prs.slide_layouts[1]
        slide = prs.slides.add_slide(slide_layout)
    
        title = slide.shapes.title
        title.text = slide_data["title"]
        title_para = title.text_frame.paragraphs[0]
        title_para.font.size = Pt(28)
        title_para.font.bold = True
        title.top = Cm(-0.5)
        title.left = Cm(0)

        remove_body(slide)
        body = slide.shapes.add_textbox(Cm(0.5), Cm(2), Inches(10.8), Inches(3.02))
        text_frame = body.text_frame
        
        for item in slide_data["content"]:
            if (item["type"] == "paragraph"):
                if isinstance(item["content"], list):
                    for p_item in item["content"]:
                        point = remove_trailing_period(p_item["point"])
                        explanation = p_item["explanation"]
                        p = text_frame.add_paragraph()  # if text_frame.text else text_frame
                        bullets.append(p)
                        point_run = p.add_run()
                        point_run.text = f"{point}: "
                        point_run.font.bold = True
                        bold_runs.append(point_run)
                        point_run.font.size = Pt(20)
                        explanation_run = p.add_run()
                        explanation_run.text = explanation
                        explanation_run.font.size = Pt(20)
                else:
                    point = remove_trailing_period(item["content"]["point"])
                    explanation = item["content"]["explanation"]
                    p = text_frame.add_paragraph()  # if text_frame.text else text_frame
                    bullets.append(p)
                    point_run = p.add_run()
                    point_run.text = f"{point}: "
                    point_run.font.bold = True
                    bold_runs.append(point_run)
                    point_run.font.size = Pt(20)
                    explanation_run = p.add_run()
                    explanation_run.text = explanation
                    explanation_run.font.size = Pt(20)
            else:
                for list_item in item["content"]:
                    if isinstance(list_item, str):
                        p = text_frame.add_paragraph()# if text_frame.text else text_frame
                        bullets.append(p)
                        p_run = p.add_run()
                        p_run.text = list_item
                        p_run.font.size = Pt(20)
                    else:
                        point = remove_trailing_period(list_item["point"])
                        explanation = list_item["explanation"]
                        p = text_frame.add_paragraph()  # if text_frame.text else text_frame
                        bullets.append(p)
                        point_run = p.add_run()
                        point_run.text = f"{point}: "
                        point_run.font.bold = True
                        bold_runs.append(point_run)
                        point_run.font.size = Pt(20)
                        explanation_run = p.add_run()
                        explanation_run.text = explanation
            text_frame.fit_text(font_family="Roboto Regular", font_file=font_path)

                        
    references_slide_layout = prs.slide_layouts[1]
    references_slide = prs.slides.add_slide(references_slide_layout)
    references_title = references_slide.shapes.title
    references_title.text = "References"
    references_title_para = references_title.text_frame.paragraphs[0]
    references_title_para.font.size = Pt(28)
    references_title_para.font.bold = True
    references_title.top = Cm(-0.5)
    references_title.left = Cm(0)

    
    remove_body(references_slide)
    references_body = references_slide.shapes.add_textbox(Cm(0.5), Cm(2), Inches(10.8), Inches(3.02))
    references_text_frame = references_body.text_frame

    for point in data["references"]:
        p = references_text_frame.add_paragraph()# if text_frame.text else text_frame
        bullets.append(p)
        p_run = p.add_run()
        p_run.text = point
        p_run.font.size = Pt(20)
    references_text_frame.fit_text(font_family=font_family, font_file=font_path)

    # make all those point runs bold
    for run in bold_runs:
      run.font.bold = True

      ## the bullets dont seem to be working
    # make level 1 bullets
    for bullet in bullets:
        bullet.level = 1

    prs.save(file_path)

def get_absolute_path(relative_path):
    """
    Converts a relative file path to an absolute file path.

    Parameters:
        relative_path (str): The relative path to be converted.

    Returns:
        str: The absolute path.
    """
    try:
        absolute_path = os.path.abspath(relative_path)
        return absolute_path
    except Exception as e:
        print(f"Error converting path: {e}")
        return None


def test_make_pptx():
    make_pptx(get_absolute_path( "./common-lisp/ninx/apps/decklm/p.pptx"), "{ \"outline\": [\"a\", \"b\", \"c\"], \"references\": [\"a\", \"b\", \"c\"], \"overall-title\":\"Testinglkkkkkkkkkkkkkkkkkkkkkkkkkkkkkkkkkkkkkkkkkkkkkkkkkkkkkkkkkkkkkkkkkkkkkkkkkkkkkkkkk\", \"content\":[{\"title\": \"Testing slide\",\"content\": [{\"type\": \"paragraph\", \"content\": {\"point\": \"test point\", \"explanation\": \"Test explanation\"}}, {\"type\": \"list\", \"content\": [\"b\", \"c\", \"d\"]}]}], \"definitions\": [{\"term\": \"test\", \"definition\": \"definition\"}]}")
